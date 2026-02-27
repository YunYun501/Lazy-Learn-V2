import pytest
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from app.services.pptx_parser import PPTXParser


def create_test_pptx(filepath: str, slides_data: list[str]):
    """Create a test PPTX with given slide texts."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for text in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        tf = txBox.text_frame
        tf.text = text

    prs.save(filepath)


def test_pptx_preserves_slide_numbers():
    """Test that PPTX parsing preserves slide numbers."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp_path = f.name

    slide_texts = ["Introduction to Z-Transform", "Properties of Z-Transform", "Applications"]
    create_test_pptx(tmp_path, slide_texts)

    parser = PPTXParser()
    slides = parser.parse(tmp_path)

    assert len(slides) == 3
    assert slides[0].slide_number == 1
    assert slides[1].slide_number == 2
    assert slides[2].slide_number == 3
    assert "Z-Transform" in slides[0].text
    assert "Properties" in slides[1].text


def test_pptx_extracts_all_text():
    """Test that all text from shapes is extracted."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp_path = f.name

    create_test_pptx(tmp_path, ["Slide with content about Z-transform and stability"])
    parser = PPTXParser()
    slides = parser.parse(tmp_path)

    assert len(slides) == 1
    assert "Z-transform" in slides[0].text


def test_pptx_to_chapters_format():
    """Test that to_chapters returns correct format."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp_path = f.name

    create_test_pptx(tmp_path, ["Slide 1 content", "Slide 2 content"])
    parser = PPTXParser()
    slides = parser.parse(tmp_path)
    chapters = parser.to_chapters(slides)

    assert len(chapters) == 2
    assert chapters[0]["number"] == "1"
    assert chapters[0]["title"] == "Slide 1"
    assert "Slide 1 content" in chapters[0]["text"]
