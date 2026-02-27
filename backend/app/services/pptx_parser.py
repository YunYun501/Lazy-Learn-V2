import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


class SlideContent:
    def __init__(self, slide_number: int, text: str, image_paths: list[str] = None):
        self.slide_number = slide_number
        self.text = text
        self.image_paths = image_paths or []


class PPTXParser:
    def __init__(self, output_dir: Path = None):
        self.output_dir = output_dir

    def parse(self, filepath: str) -> list[SlideContent]:
        """Parse a PPTX file. Returns list of SlideContent with slide numbers."""
        prs = Presentation(filepath)
        slides = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Extract all text from shapes
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

            slide_text = "\n".join(text_parts)

            # Extract images from shapes
            image_paths = []
            if self.output_dir:
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            image = shape.image
                            ext = image.ext
                            img_path = self.output_dir / f"slide{slide_num}_img{shape_idx}.{ext}"
                            img_path.write_bytes(image.blob)
                            image_paths.append(str(img_path))
                        except Exception:
                            pass

            slides.append(SlideContent(
                slide_number=slide_num,
                text=slide_text,
                image_paths=image_paths,
            ))

        return slides

    def to_chapters(self, slides: list[SlideContent]) -> list[dict]:
        """Convert slides to chapter-like structure for unified processing."""
        return [
            {
                "number": str(s.slide_number),
                "title": f"Slide {s.slide_number}",
                "text": s.text,
                "page_start": s.slide_number,
                "page_end": s.slide_number,
                "image_paths": s.image_paths,
            }
            for s in slides
        ]
